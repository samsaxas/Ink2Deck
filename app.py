import streamlit as st
import base64
import os
from pymongo import MongoClient
import hashlib
import certifi
from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from fpdf import FPDF
import numpy as np
import cv2
import pytesseract
import google.generativeai as genai

# Initialize session state for navigation
if "page" not in st.session_state:
    st.session_state.page = "home"
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "username" not in st.session_state:
    st.session_state.username = ""

# Set page config with black background
st.set_page_config(
    page_title="Ink2Deck",
    page_icon="📝",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Custom CSS for full black background and button styling
st.markdown(
    """
    <style>
    .stApp {
        background-color: #000000;
        color: white;
    }
    .stTextInput>div>div>input,
    .stTextArea>div>div>textarea,
    .stSelectbox>div>div>select,
    .stButton>button {
        background-color: #111111;
        color: white;
        border: 1px solid #077b32;
    }
    .stTextInput>label,
    .stTextArea>label,
    .stSelectbox>label,
    .stFileUploader>label {
        color: white !important;
    }
    .st-bb {
        background-color: #000000;
    }
    .st-at {
        background-color: #111111;
    }
    .st-cj {
        background-color: #077b32;
    }
    .st-ck {
        background-color: #077b32;
    }
    .st-cl {
        background-color: #077b32;
    }
    .st-cm {
        background-color: #077b32;
    }
    .st-cn {
        background-color: #077b32;
    }
    .stTab {
        background-color: #000000;
    }
    .stTabs [data-baseweb="tab-list"] {
        background-color: #000000;
    }
    .stTabs [aria-selected="true"] {
        background-color: #077b32 !important;
        color: white !important;
    }
    .stTabs [aria-selected="false"] {
        background-color: #111111 !important;
        color: white !important;
    }
    /* Custom style for the Get Started button */
    div.stButton > button:first-child {
        width: 50%;
        height: 6vh;
        margin: 30px auto;
        background-color: #077b32;
        color: white;
        border: none;
        outline: none;
        font-size: 120%;
        font-weight: 700;
        border-radius: 5px;
        transition: all 0.2s linear;
        display: block;
    }
    div.stButton > button:first-child:hover {
        transform: scale(1.1);
        color: #077b32;
        border: 2px solid #077b32;
        background-color: transparent;
        font-weight: 700;
        box-shadow: 0 0 40px #077b32;
    }
    /* Chatbot styles - Only on home page */
    .orimon-chatbot {
        position: fixed;
        bottom: 20px;
        right: 20px;
        z-index: 9999;
    }
    .orimon-chat-icon {
        width: 60px;
        height: 60px;
        border-radius: 50%;
        background-color: #077b32;
        display: flex;
        justify-content: center;
        align-items: center;
        cursor: pointer;
        box-shadow: 0 0 20px rgba(7, 123, 50, 0.5);
        transition: all 0.3s ease;
    }
    .orimon-chat-icon:hover {
        transform: scale(1.1);
        box-shadow: 0 0 30px rgba(7, 123, 50, 0.8);
    }
    .orimon-chat-icon i {
        color: white;
        font-size: 24px;
    }
    .orimon-chat-window {
        width: 350px;
        height: 500px;
        border-radius: 10px;
        overflow: hidden;
        box-shadow: 0 0 20px rgba(0, 0, 0, 0.5);
        display: none;
    }
    @media (max-width: 884px) {
        div.stButton > button:first-child {
            width: 80%;
        }
        .orimon-chat-window {
            width: 300px;
            height: 450px;
        }
    }
    @media (max-width: 440px) {
        div.stButton > button:first-child {
            width: 90%;
            margin: 20px auto;
        }
        .orimon-chat-window {
            width: 280px;
            height: 400px;
            right: 10px;
        }
    }
    </style>
    """,
    unsafe_allow_html=True
)

# =============================================
# Common Functions
# =============================================

def get_image_base64(path):
    try:
        normalized_path = os.path.normpath(path)
        if not os.path.exists(normalized_path):
            return None
        with open(normalized_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode()
    except Exception:
        return None

def load_image():
    possible_paths = [
        "main2.jpg", "main4.jpg",
        os.path.join("main2.jpg"),
        os.path.join("main4.jpg"),
        os.path.join(os.path.dirname(__file__), "main2.jpg"),
        os.path.join(os.path.dirname(__file__), "Ink2Deck", "main2.jpg"),
        "/app/Ink2Deck/main2.jpg", "/app/main2.jpg"
    ]
    for img_path in possible_paths:
        img_base64 = get_image_base64(img_path)
        if img_base64:
            return img_base64
    return None

# =============================================
# Home Page (with chatbot)
# =============================================

def home_page():
    # Get the base64 encoded image
    img_data = load_image() or ""
    
    # Create the HTML content
    html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.5.1/css/all.min.css">
    <title>Ink2Deck</title>
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@100;200;300;400;500;600&display=swap');
        *{{
            padding: 0;
            margin: 0;
            font-family: 'Poppins', sans-serif;
            box-sizing: border-box;
        }}

        body{{
            width: 100%;
            height: 100vh;
            overflow: hidden;
            background-color: black;
        }}
        nav{{
            width: 100%;
            height: 10vh;
            position: sticky;
        }}
        .nav-container{{
            width: 100%;
            height: 100%;
            display: flex;
            justify-content: space-around;
            align-items: center;
        }}
        .logo{{
            color: white;
            font-size: 2rem;
            font-weight: bold;
        }}
        .logo span{{
            color: #077b32;
            text-shadow: 0 0 10px #077b32;
        }}
        .main-container{{
            width: 100%;
            height: 90vh;
            display: flex;
            justify-content: space-evenly;
            align-items: center;
        }}
        .main-container .image-container {{
            width: 500px;
            height: 500px;
            border-radius: 50%;
            overflow: hidden;
            box-shadow: 0 0 50px #077b32;
            display: flex;
            justify-content: center;
            align-items: center;
        }}
        .main-container .image-container img{{
            width: 100%;
            height: 100%;
            object-fit: cover;
        }}
        .main-container .content{{
            color: white;
            width: 40%;
        }}
        .content h1{{
            font-size: clamp(1rem, 1rem + 5vw, 1.8rem);
        }}
        .content h1 span{{
            color: #077b32;
            text-shadow: 0  0 10px #077b32;
        }}
        .content .deck-maker{{
            font-size: clamp(1rem, 1rem + 5vw, 2.5rem);
            font-weight: 600;
        }}
        .content .deck-maker span{{
            color: #077b32;
            text-shadow: 0 0 10px #077b32;
        }}
        .content p{{
            font-size: clamp(0.4rem , 0.2rem + 9vw, 1rem);
            margin: 10px 0 30px 0;
        }}
        @media (max-width:884px) {{
            .main-container {{
                flex-direction: column;
            }}
            .main-container .content{{
                width: 80%;
            }}
            .main-container .image-container{{
                width: 300px;
                height: 300px;
            }}
        }}
        @media (max-width:440px){{
            .main-container .image-container{{
                width: 250px;
                height: 250px;
            }}
        }}
    </style>
</head>
<body>
    <nav>
        <div class="nav-container">
            <div class="logo">
                Ink2<span>Deck</span>
            </div>
        </div>
    </nav>
    <div class="main-container">
        <div class="image-container">
            <img src="data:image/png;base64,{img_data}" alt="Ink2Deck">
        </div>
        <div class="content">
            <h1>Hey I'm <span>Ink2Deck</span></h1>
            <div class="deck-maker">I'm a <span>DECK MAKER</span></div>
            <p>The Ink2Deck captures whiteboard content and converts handwritten notes, diagrams, and equations into structured slide decks. It generates downloadable PowerPoint or PDF files, ensuring clarity and easy sharing.</p>
        </div>
    </div>
</body>
</html>
"""
    st.markdown(html_content, unsafe_allow_html=True)
    
    # Add the Streamlit button
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        if st.button("Get Started!", key="get_started"):
            st.session_state.page = "login"
            st.rerun()
    
    # Add Orimon AI Chatbot (only on home page)
    st.markdown(
        """
        <div class="orimon-chatbot">
            <div class="orimon-chat-icon" id="chat-icon">
                <i class="fas fa-comment-dots"></i>
            </div>
            <div class="orimon-chat-window" id="chat-window"></div>
        </div>
        """,
        unsafe_allow_html=True
    )

# =============================================
# Login/Signup Page
# =============================================

def login_page():
    # Custom CSS for the login container
    st.markdown(
        """
        <style>
            /* Main container styling */
            .login-box {
                background-color: #ffffff;
                border-radius: 10px;
                padding: 2rem;
                width: 350px;
                margin: 2rem auto;
                box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
            }
            
            /* Input field styling */
            .stTextInput>div>div>input {
                background-color: #f8f9fa !important;
                color: #333 !important;
                border: 1px solid #ced4da !important;
                width: 100% !important;
            }
            
            /* Button styling */
            .stButton>button {
                width: 100% !important;
                margin: 0.5rem 0 !important;
                transition: all 0.3s ease !important;
            }
            
            /* Primary button */
            div[data-testid="stButton"]:has(button[kind="primary"]) button {
                background-color: #077b32 !important;
                color: white !important;
                border: none !important;
            }
            
            /* Secondary button */
            div[data-testid="stButton"]:has(button[kind="secondary"]) button {
                background-color: transparent !important;
                color: #077b32 !important;
                border: 1px solid #077b32 !important;
            }
            
            /* Tab styling */
            [data-baseweb="tab-list"] {
                gap: 0.5rem !important;
            }
            
            button[data-baseweb="tab"] {
                background-color: #f8f9fa !important;
                color: #333 !important;
                border-radius: 4px !important;
                padding: 0.5rem 1rem !important;
                flex: 1 !important;
            }
            
            button[data-baseweb="tab"][aria-selected="true"] {
                background-color: #077b32 !important;
                color: white !important;
            }
        </style>
        """,
        unsafe_allow_html=True
    )

    # MongoDB connection (keep your existing connection code)
    load_dotenv()
    try:
        MONGO_URI = os.getenv("MONGO_URI")
        if not MONGO_URI:
            raise ValueError("MongoDB URI not found")
        client = MongoClient(
            MONGO_URI,
            tlsCAFile=certifi.where(),
            serverSelectionTimeoutMS=5000,
            connectTimeoutMS=3000
        )
        db = client["Ink2Deck"]
        users = db["users"]
    except Exception as e:
        st.error(f"Database connection failed: {str(e)}")
        return

    # Create the white box container
    with st.container():
        st.markdown('<div class="login-box">', unsafe_allow_html=True)
        
        # Create tabs for Login and Sign Up
        tab1, tab2 = st.tabs(["Login", "Sign Up"])
        
        with tab1:
            with st.form("login_form"):
                st.header("Login")
                username = st.text_input("Username", placeholder="Enter your username")
                password = st.text_input("Password", type="password", placeholder="••••••••")
                
                if st.form_submit_button("Login", type="primary"):
                    user = users.find_one({"username": username})
                    if user and user["password"] == hashlib.sha256(password.encode()).hexdigest():
                        st.session_state.logged_in = True
                        st.session_state.username = username
                        st.session_state.page = "upload"
                        st.rerun()
                    else:
                        st.error("Invalid username or password")
        
        with tab2:
            with st.form("signup_form"):
                st.header("Create Account")
                name = st.text_input("Full Name", placeholder="Your full name")
                email = st.text_input("Email", placeholder="your@email.com")
                username = st.text_input("Username", placeholder="Choose a username")
                password = st.text_input("Password", type="password", placeholder="••••••••")
                terms = st.checkbox("I agree to the Terms and Conditions")
                
                if st.form_submit_button("Create Account", type="primary"):
                    if not terms:
                        st.error("Please accept the Terms and Conditions")
                    elif not all([name, email, username, password]):
                        st.error("Please fill all fields")
                    elif users.find_one({"$or": [{"username": username}, {"email": email}]}):
                        st.error("Username or email already exists")
                    else:
                        users.insert_one({
                            "name": name,
                            "email": email,
                            "username": username,
                            "password": hashlib.sha256(password.encode()).hexdigest()
                        })
                        st.success("Account created! Please login.")
        
        if st.button("Back to Home", type="secondary"):
            st.session_state.page = "home"
            st.rerun()
        
        st.markdown('</div>', unsafe_allow_html=True)

# =============================================
# Upload/Processing Page
# =============================================

def upload_page():
    # Initialize Gemini
    genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
    model = genai.GenerativeModel('gemini-1.5-pro')
    
    st.title("Upload Whiteboard Image")
    st.write(f"Welcome, {st.session_state.username}!")
    
    if not st.session_state.get("logged_in", False):
        st.warning("Please login first")
        st.session_state.page = "login"
        st.rerun()
        return
    
    def preprocess_image(image):
        """Optimized preprocessing pipeline"""
        img_array = np.array(image)
        gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)
        return Image.fromarray(thresh)

    def extract_text_with_gemini(image):
        """Use Gemini Vision for superior handwriting extraction"""
        try:
            img_byte_arr = BytesIO()
            image.save(img_byte_arr, format='PNG')
            img_bytes = img_byte_arr.getvalue()

            response = model.generate_content([
                "Extract all text from this whiteboard/image exactly as written, including equations. "
                "Preserve line breaks and original language.",
                {"mime_type": "image/png", "data": img_bytes}
            ])
            
            return response.text if hasattr(response, 'text') else ""
        except Exception as e:
            st.warning(f"Gemini extraction failed: {e}")
            return ""

    def extract_text_with_tesseract(image):
        """Extract text using Tesseract OCR with preprocessing"""
        try:
            # Convert PIL Image to OpenCV format
            img_array = np.array(image)
            
            # Convert to grayscale
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
            
            # Apply thresholding to preprocess the image
            thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY | cv2.THRESH_OTSU)[1]
            
            # Convert back to PIL for Tesseract
            thresh_image = Image.fromarray(thresh)
            
            return pytesseract.image_to_string(thresh_image, config='--psm 6')
        except Exception as e:
            st.warning(f"Tesseract extraction failed: {e}")
            return ""

    def extract_text_with_ocr(image):
        """Optimized text extraction with fallback logic"""
        try:
            # Try Gemini first if API key exists
            if os.getenv("GEMINI_API_KEY"):
                gemini_text = extract_text_with_gemini(image)
                if gemini_text.strip():
                    return gemini_text
            
            # Fallback to Tesseract
            return extract_text_with_tesseract(image)
        except Exception:
            return ""

    def create_ppt_with_text_and_image(text, image):
        """Create PowerPoint presentation"""
        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "Whiteboard Content"
        title_slide.placeholders[1].text = "Automatically Generated from Image"
        
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        image_stream = BytesIO()
        image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, Inches(1), Inches(1), width=Inches(6))

        slides_content = [s.strip() for s in text.split("\n\n") if s.strip()]
        for i, content in enumerate(slides_content, 1):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Slide {i}"
            slide.placeholders[1].text = content
            for paragraph in slide.placeholders[1].text_frame.paragraphs:
                paragraph.font.size = Pt(18)

        ppt_stream = BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        return ppt_stream

    def create_pdf(text):
        """Create PDF with better Unicode support"""
        pdf = FPDF()
        pdf.add_page()
        try:
            pdf.add_font('DejaVu', '', 'DejaVuSansCondensed.ttf', uni=True)
            pdf.set_font('DejaVu', '', 12)
        except:
            try:
                pdf.add_font('Arial', '', 'arial.ttf', uni=True)
                pdf.set_font('Arial', '', 12)
            except:
                pdf.set_font('Arial', '', 12)
        
        try:
            pdf.multi_cell(0, 10, txt=text)
        except:
            cleaned_text = text.encode('utf-8', errors='ignore').decode('utf-8')
            pdf.multi_cell(0, 10, txt=cleaned_text)
        
        return pdf.output(dest='S').encode('latin-1', errors='replace')

    # File uploader
    uploaded_file = st.file_uploader("Choose an image (JPG, PNG)", type=["png", "jpg", "jpeg"])
    
    if uploaded_file is not None:
        try:
            image = Image.open(uploaded_file)
            
            with st.spinner("Processing image and extracting text..."):
                processed_image = preprocess_image(image)
                extracted_text = extract_text_with_ocr(processed_image)
            
            if extracted_text.strip():
                # Create two columns for image and text
                col1, col2 = st.columns(2)
                
                with col1:
                    st.image(image, caption="Uploaded Image", use_container_width=True)
                
                with col2:
                    st.success("Text extracted successfully!")
                    st.text_area("Extracted Text", extracted_text, height=400)
                
                # Download and navigation buttons
                st.write("---")  # Add a divider
                
                pdf_bytes = create_pdf(extracted_text)
                ppt_stream = create_ppt_with_text_and_image(extracted_text, image)
                
                # Center-align the buttons
                col1, col2, col3, col4, col5 = st.columns([1,2,2,2,1])
                
                with col2:
                    logout_button = st.button("Logout")
                    if logout_button:
                        st.session_state.logged_in = False
                        st.session_state.page = "login"
                        st.rerun()
                
                with col3:
                    st.download_button(
                        label="Download PDF",
                        data=pdf_bytes,
                        file_name="extracted_content.pdf",
                        mime="application/pdf"
                    )
                
                with col4:
                    st.download_button(
                        label="Download PowerPoint",
                        data=ppt_stream,
                        file_name="whiteboard_presentation.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                
                # Centered Back to Home button
                st.write("---")
                col1, col2, col3 = st.columns([3,2,3])
                with col2:
                    if st.button("Back to Home"):
                        st.session_state.page = "home"
                        st.rerun()
            else:
                st.error("No text detected. Please try a clearer image.")
                
        except Exception as e:
            st.error(f"An error occurred: {str(e)}")
    else:
        # Add back and logout buttons when no file is uploaded
        col1, col2, col3 = st.columns([3,2,3])
        with col1:
            if st.button("Logout"):
                st.session_state.logged_in = False
                st.session_state.page = "login"
                st.rerun()
        with col3:
            if st.button("Back to Home"):
                st.session_state.page = "home"
                st.rerun()

                
# =============================================
# Main App Router
# =============================================

if st.session_state.page == "home":
    home_page()
elif st.session_state.page == "login":
    login_page()
elif st.session_state.page == "upload":
    upload_page()

# Add Orimon AI Chatbot script (only loads on home page)
st.markdown(
    """
    <script>
    // Only load chatbot on home page
    if (window.location.pathname === '/' || window.location.pathname.endsWith('Ink2Deck/')) {
        // Load chatbot script
        const script = document.createElement('script');
        script.src = 'https://bot.orimon.ai/deploy/index.js';
        script.setAttribute('tenantId', 'aa43bce8-3d59-46ea-8320-2ea2900b55de');
        document.body.appendChild(script);
        
        // Add toggle functionality
        document.getElementById('chat-icon').addEventListener('click', function() {
            const chatWindow = document.getElementById('chat-window');
            if (chatWindow.style.display === 'block') {
                chatWindow.style.display = 'none';
            } else {
                chatWindow.style.display = 'block';
            }
        });
    }
    </script>
    """,
    unsafe_allow_html=True
)